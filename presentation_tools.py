"""
presentation_tools.py
---------------------
json_to_ppt(
    data: Union[Dict[str, Any], List[Dict[str, Any]]],
    address: str,
    lat: float,
    lon: float
) -> str

Produces a 3-slide PowerPoint:
  1) Centered header slide (big text)
  2) Table slide from `data`
  3) Two images side-by-side: a Google Static Map & a Street View photo
"""

import os
import tempfile
import uuid
import textwrap
from typing import Any, Dict, List, Union
from pathlib import Path

import requests
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# RTL shaping (optional)
try:
    import arabic_reshaper
    from bidi.algorithm import get_display
    _RTL = True
except ImportError:
    _RTL = False

GOOGLE_API_KEY = "AIzaSyDVF3vo_gSNIQvenS8FSSRu_gb2f-kM6bs"


def _rtl_header(text: str, wrap: int = 15) -> str:
    wrapped = "\n".join(textwrap.wrap(text, wrap))
    if _RTL:
        reshaped = arabic_reshaper.reshape(wrapped)
        return get_display(reshaped)
    return "\n".join(line[::-1] for line in wrapped.split("\n"))


def _table_image(rows: List[Dict[str, Any]]) -> Path:
    headers = [_rtl_header(h) for h in rows[0].keys()]
    cells   = [[row.get(h, "") for h in rows[0].keys()] for row in rows]

    fig_w = max(6, len(headers) * 0.9)
    fig_h = max(1.5, len(rows) * 0.6)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.axis("off")

    tbl = ax.table(
        cellText=cells,
        colLabels=headers,
        loc="center",
        cellLoc="center",
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(10)
    tbl.auto_set_column_width(col=list(range(len(headers))))

    for (r, c), cell in tbl._cells.items():
        cell.get_text().set_ha("right")
        cell.get_text().set_va("center")
        if r == 0:
            cell.set_facecolor("#1976D2")
            cell.get_text().set_color("white")
            cell.set_height(cell.get_height() * 2)

    out = Path(tempfile.gettempdir()) / f"table_{uuid.uuid4().hex}.png"
    plt.savefig(out, bbox_inches="tight", dpi=150)
    plt.close(fig)
    return out


def _fetch_image(url: str, timeout: int = 10) -> Path:
    resp = requests.get(url, timeout=timeout)
    resp.raise_for_status()
    out = Path(tempfile.gettempdir()) / f"img_{uuid.uuid4().hex}.png"
    out.write_bytes(resp.content)
    return out


def _set_slide_bg(slide, color_hex: str):
    """Fill slide background with solid color (hex without '#')."""
    fill = slide.background.fill
    fill.solid()
    r = int(color_hex[0:2], 16)
    g = int(color_hex[2:4], 16)
    b = int(color_hex[4:6], 16)
    fill.fore_color.rgb = RGBColor(r, g, b)


def json_to_ppt(
    data: Union[Dict[str, Any], List[Dict[str, Any]]],
    address: str,
    lat: float,
    lon: float
) -> str:
    # Normalize data → list of rows
    if isinstance(data, dict) and "data" in data:
        rows = data["data"]
    elif isinstance(data, list):
        rows = data
    elif isinstance(data, dict):
        rows = [data]
    else:
        raise ValueError(f"Unsupported data format: {type(data)}")

    prs = Presentation()
    SLIDE_BG = "BBDEFB"  # light-blue

    # Slide 1: centered big header
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide1, SLIDE_BG)
    sw, sh = prs.slide_width, prs.slide_height
    box_w, box_h = Inches(8), Inches(2)
    left = (sw - box_w) / 2
    top  = (sh - box_h) / 2
    tb = slide1.shapes.add_textbox(left, top, box_w, box_h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = address
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True

    # Slide 2: table
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide2, SLIDE_BG)
    table_png = _table_image(rows)
    slide2.shapes.add_picture(
        str(table_png),
        left=Inches(0.5), top=Inches(0.5),
        width=Inches(9)
    )

    # Slide 3: two Google images
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide3, SLIDE_BG)

    # Static map
    map_img = None
    if GOOGLE_API_KEY:
        map_url = (
            "https://maps.googleapis.com/maps/api/staticmap"
            f"?center={lat},{lon}&zoom=18&size=400x300"
            f"&markers=color:red|{lat},{lon}"
            f"&key={GOOGLE_API_KEY}"
        )
        try:
            map_img = _fetch_image(map_url)
        except Exception:
            map_img = None

    # Street view
    sv_img = None
    if GOOGLE_API_KEY:
        sv_url = (
            "https://maps.googleapis.com/maps/api/streetview"
            f"?size=400x300&location={lat},{lon}"
            f"&fov=90&heading=235&pitch=10"
            f"&key={GOOGLE_API_KEY}"
        )
        try:
            sv_img = _fetch_image(sv_url)
        except Exception:
            sv_img = None

    left = Inches(0.5)
    if map_img:
        slide3.shapes.add_picture(str(map_img), left, Inches(0.75), width=Inches(4))
    else:
        ph = slide3.shapes.add_textbox(left, Inches(2), Inches(4), Inches(0.5)).text_frame
        ph.text = "מפת מיקום לא זמינה"
        ph.paragraphs[0].font.size = Pt(18)
        ph.paragraphs[0].alignment = PP_ALIGN.CENTER

    if sv_img:
        slide3.shapes.add_picture(str(sv_img), left + Inches(4.5), Inches(0.75), width=Inches(4))
    else:
        ph2 = slide3.shapes.add_textbox(left + Inches(4.5), Inches(2), Inches(4), Inches(0.5)).text_frame
        ph2.text = "תצפית רחוב לא זמינה"
        ph2.paragraphs[0].font.size = Pt(18)
        ph2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save PPTX
    out_ppt = Path.cwd() / f"PlanTable_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_ppt)
    return str(out_ppt.resolve())


# Expose for AutoGen
TOOL_MAP = {"json_to_ppt": json_to_ppt}
